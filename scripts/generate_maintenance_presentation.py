"""
Generates a professional PowerPoint summarizing continuous development & maintenance work.

Run: python backend/scripts/generate_maintenance_presentation.py

Dependencies: python-pptx, matplotlib, numpy
The script uses example/sample data; update the `DATA` dict to reflect real project numbers and items.
"""
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import numpy as np
import os

OUT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
ASSETS_DIR = os.path.join(OUT_DIR, 'presentation_assets')
if not os.path.exists(ASSETS_DIR):
    os.makedirs(ASSETS_DIR)

OUTPUT_PPTX = os.path.join(OUT_DIR, 'maintenance_report.pptx')

# ----- SAMPLE DATA (edit this section with real values) -----
DATA = {
    'app_name': 'Salon Management Platform',
    'duration': 'Last 6 months (Nov 2025 - Apr 2026)',
    'team': 'Platform & Backend Team',
    'timeline_months': ['Nov', 'Dec', 'Jan', 'Feb', 'Mar', 'Apr'],
    # stacked categories per month: features, bug_fixes, perf, infra
    'timeline_values': np.array([
        [2, 5, 3, 1],
        [3, 4, 2, 2],
        [4, 3, 5, 1],
        [2, 6, 4, 2],
        [5, 2, 6, 3],
        [3, 3, 4, 2],
    ]),
    'git_heatmap': np.random.randint(0, 6, size=(7, 26)).tolist(),  # weeks x days
    'metrics': {
        'commits': 482,
        'deployments': 24,
        'bugs_resolved': 138,
        'features_shipped': 27,
        'migrations': 6,
        'api_improvements': 18,
        'uptime_percent': 99.97,
    },
    'features': [
        ('Automated Booking Retry', 'Reduces failed bookings during transient errors', 'Improved reliability; fewer support tickets'),
        ('Search Index Optimization', 'Faster product/service search', 'Better user experience; lower DB load'),
        ('Payments Retry & Idempotency', 'Prevents double-charges on network hiccups', 'Reduced payment disputes'),
    ],
    'backend_changes': [
        'DB schema changes and safe data migrations',
        'API contract versioning and deprecation handling',
        'Query tuning and index additions (CONCURRENTLY where possible)',
        'Enhanced logging, monitoring, and alerts',
        'Auth hardening and token rotation automation',
    ],
}

# ----- Helper image/chart generators -----

def save_timeline_chart(months, values, path):
    # values shape: (n_months, 4)
    categories = ['Features', 'Bug Fixes', 'Performance', 'Infrastructure']
    colors = ['#2F6DB3', '#6EA4D8', '#8FB7D9', '#9AAFCB']
    values = np.array(values)
    fig, ax = plt.subplots(figsize=(10, 2.4))
    bottom = np.zeros(len(months))
    for i in range(values.shape[1]):
        ax.bar(months, values[:, i], bottom=bottom, color=colors[i], label=categories[i])
        bottom += values[:, i]
    ax.set_ylabel('Items / Month')
    ax.set_title('Timeline: Monthly Activity (stacked)')
    ax.legend(loc='upper left', bbox_to_anchor=(1, 1))
    plt.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)


def save_git_heatmap(heatmap, path):
    data = np.array(heatmap)
    # transpose so weeks across x-axis
    fig, ax = plt.subplots(figsize=(10, 1.8))
    c = ax.imshow(data.T, cmap='YlGn', aspect='auto', interpolation='nearest')
    ax.set_yticks([0, data.shape[1]-1])
    ax.set_ylabel('Days')
    ax.set_xlabel('Weeks')
    ax.set_title('Git Activity Heatmap (sample)')
    plt.colorbar(c, orientation='vertical', fraction=0.02, pad=0.02)
    plt.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)


def save_metrics_chart(metrics, path):
    keys = ['commits', 'deployments', 'bugs_resolved', 'features_shipped', 'migrations', 'api_improvements']
    labels = ['Commits', 'Deployments', 'Bugs Resolved', 'Features', 'Migrations', 'API improv.']
    values = [metrics[k] for k in keys]
    fig, ax = plt.subplots(figsize=(8, 3))
    bars = ax.bar(labels, values, color='#2F6DB3')
    ax.set_title('Key Metrics')
    ax.set_ylabel('Count')
    plt.xticks(rotation=30, ha='right')
    for bar, v in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2, v + max(values)*0.01, str(v), ha='center', va='bottom', fontsize=9)
    plt.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)


def save_architecture_diagram(path):
    # Simple layered boxes diagram
    fig, ax = plt.subplots(figsize=(6, 4))
    ax.axis('off')
    colors = ['#DCEAF6', '#B9D3F1', '#96BBEE', '#6EA4D8']
    labels = ['Presentation (UI)', 'API Layer', 'Business Logic / Services', 'Data / DB & Infra']
    y = 0.8
    for i, (c, lbl) in enumerate(zip(colors[::-1], labels[::-1])):
        rect = plt.Rectangle((0.1, y - 0.18), 0.8, 0.15, facecolor=c, edgecolor='black')
        ax.add_patch(rect)
        ax.text(0.5, y - 0.105, lbl, ha='center', va='center', fontsize=12)
        y -= 0.18
    ax.set_title('Simplified System Layers')
    plt.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)


# ----- PowerPoint generation -----

def add_title_slide(prs, title, subtitle, author):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    title_tf.text = title
    subtitle_tf = slide.placeholders[1].text_frame
    subtitle_tf.text = subtitle + '\n\n' + author


def add_bullets_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(14)


def add_image_slide(prs, title, image_path, caption=None):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.4)
    pic = slide.shapes.add_picture(image_path, left, top, width=Inches(9))
    if caption:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(9), Inches(0.6))
        tf = txBox.text_frame
        tf.text = caption
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.italic = True


def build_presentation(data):
    prs = Presentation()
    # set slide size to widescreen (optional)
    from pptx.util import Cm
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    # Title
    add_title_slide(prs, f"{data['app_name']} — Continuous Development & Maintenance",
                    data['duration'], f"Prepared by {data['team']} — {datetime.now().strftime('%b %d, %Y')}")

    # Project Overview
    bullets = [
        f"Application: {data['app_name']}",
        f"Maintenance period: {data['duration']}",
        f"Team: {data['team']}",
        "Scope: Backend, infrastructure, database, API, performance, security, feature work, monitoring",
    ]
    add_bullets_slide(prs, 'Project Overview', bullets)

    # Timeline: create chart image and add slide
    timeline_img = os.path.join(ASSETS_DIR, 'timeline.png')
    save_timeline_chart(data['timeline_months'], data['timeline_values'], timeline_img)
    add_image_slide(prs, 'Continuous Development Timeline', timeline_img,
                    'Stacked monthly activity: features, bug fixes, performance, infrastructure')

    # Git activity
    heatmap_img = os.path.join(ASSETS_DIR, 'git_heatmap.png')
    save_git_heatmap(data['git_heatmap'], heatmap_img)
    add_image_slide(prs, 'Git Activity & Contributions', heatmap_img,
                    'Each cell represents commit activity (darker = more commits). Commits include bug fixes, DB changes, refactors, and features.')

    # Backend & Database
    add_bullets_slide(prs, 'Backend & Database Work (summary)', data['backend_changes'])

    # New Features Delivered
    feature_bullets = []
    for name, problem, impact in data['features']:
        feature_bullets.append(f"{name}: {problem} — Impact: {impact}")
    add_bullets_slide(prs, 'New Features Delivered', feature_bullets)

    # Maintenance & Stability
    maint_bullets = [
        'Ongoing bug fixes and incident response',
        'Dependency and security patching',
        'Refactoring for maintainability',
        'Monitoring, alerting, and SLO improvements',
        'Scalability & capacity planning efforts',
        'Automations for deployments and rollbacks',
    ]
    add_bullets_slide(prs, 'Maintenance & Stability', maint_bullets)

    # Metrics & Impact
    metrics_img = os.path.join(ASSETS_DIR, 'metrics.png')
    save_metrics_chart(data['metrics'], metrics_img)
    add_image_slide(prs, 'Metrics & Impact', metrics_img,
                    f"High-level metrics including commits, deployments, bugs resolved, features shipped. Uptime: {data['metrics'].get('uptime_percent', 'n/a'):.2f}%")

    # Technical Complexity simplified (architecture)
    arch_img = os.path.join(ASSETS_DIR, 'architecture.png')
    save_architecture_diagram(arch_img)
    add_image_slide(prs, 'How Backend Supports the UI', arch_img,
                    'Simplified layers: Presentation → API → Services → Data & Infrastructure')

    # Conclusion
    concl_bullets = [
        'Continuous engineering reduces incidents and improves user experience over time',
        'Many improvements are backend-focused and not always visible in the UI',
        'We prioritize reliability, scalability, security, and measurable business value',
        'Next steps: continue roadmap delivery, share regular status dashboards, and measure impact',
    ]
    add_bullets_slide(prs, 'Conclusion', concl_bullets)

    # Save
    prs.save(OUTPUT_PPTX)
    print('Saved presentation to:', OUTPUT_PPTX)


if __name__ == '__main__':
    build_presentation(DATA)
