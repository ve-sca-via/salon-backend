"""
Build a professional client-facing presentation summarising development
and maintenance work delivered on the salon platform backend.

All numbers are sourced from the live git repository (commits, PRs,
migrations, services).  Run from the repo root:

    python scripts/build_client_presentation.py

Output: ./Development_Progress_Report.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------

NAVY = RGBColor(0x0F, 0x2A, 0x4A)
PRIMARY = RGBColor(0x1F, 0x4E, 0x8C)
ACCENT = RGBColor(0x2E, 0x86, 0xDE)
LIGHT = RGBColor(0xE9, 0xF1, 0xFA)
SOFT = RGBColor(0xF5, 0xF7, 0xFA)
INK = RGBColor(0x1A, 0x1F, 0x2C)
MUTED = RGBColor(0x6B, 0x73, 0x80)
LINE = RGBColor(0xCF, 0xD8, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x27, 0xAE, 0x60)
AMBER = RGBColor(0xE6, 0x7E, 0x22)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return slide


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=14,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font="Calibri",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, fill, line=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    if not shadow:
        shape.shadow.inherit = False
    return shape


def add_round_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_oval(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_line(slide, x1, y1, x2, y2, color=LINE, weight=1.0):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(weight)
    return connector


def add_page_header(slide, eyebrow, title):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.18), PRIMARY)
    add_text(
        slide,
        Inches(0.6),
        Inches(0.35),
        Inches(10),
        Inches(0.35),
        eyebrow.upper(),
        size=11,
        bold=True,
        color=ACCENT,
    )
    add_text(
        slide,
        Inches(0.6),
        Inches(0.6),
        Inches(12),
        Inches(0.7),
        title,
        size=28,
        bold=True,
        color=NAVY,
    )
    add_rect(slide, Inches(0.6), Inches(1.25), Inches(0.6), Inches(0.05), ACCENT)


def add_footer(slide, page_label):
    add_rect(slide, 0, Inches(7.32), SLIDE_W, Inches(0.18), LIGHT)
    add_text(
        slide,
        Inches(0.6),
        Inches(7.18),
        Inches(8),
        Inches(0.3),
        "Salon Platform — Development Progress Report",
        size=9,
        color=MUTED,
    )
    add_text(
        slide,
        Inches(11.5),
        Inches(7.18),
        Inches(1.5),
        Inches(0.3),
        page_label,
        size=9,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def slide_cover(prs):
    s = blank_slide(prs)

    # Left band
    add_rect(s, 0, 0, Inches(4.6), SLIDE_H, NAVY)
    add_rect(s, 0, Inches(3.4), Inches(4.6), Inches(0.06), ACCENT)

    add_text(
        s,
        Inches(0.6),
        Inches(0.6),
        Inches(4),
        Inches(0.4),
        "CLIENT REPORT",
        size=12,
        bold=True,
        color=ACCENT,
    )
    add_text(
        s,
        Inches(0.6),
        Inches(2.4),
        Inches(3.6),
        Inches(0.5),
        "Confidential",
        size=11,
        color=WHITE,
    )
    add_text(
        s,
        Inches(0.6),
        Inches(2.75),
        Inches(3.6),
        Inches(0.45),
        "Period: Nov 2025 — May 2026",
        size=12,
        bold=True,
        color=WHITE,
    )
    add_text(
        s,
        Inches(0.6),
        Inches(3.55),
        Inches(3.6),
        Inches(2),
        "Backend engineering,\ndatabase, infrastructure\nand product delivery",
        size=14,
        color=LIGHT,
    )

    # Right hero area
    add_text(
        s,
        Inches(5.2),
        Inches(1.8),
        Inches(7.7),
        Inches(0.6),
        "Development & Maintenance",
        size=24,
        color=PRIMARY,
        bold=True,
    )
    add_text(
        s,
        Inches(5.2),
        Inches(2.4),
        Inches(7.7),
        Inches(1.4),
        "Progress Report",
        size=54,
        color=NAVY,
        bold=True,
    )
    add_rect(s, Inches(5.2), Inches(4.05), Inches(1.0), Inches(0.07), ACCENT)
    add_text(
        s,
        Inches(5.2),
        Inches(4.2),
        Inches(7.7),
        Inches(1.6),
        "A transparent summary of the engineering, infrastructure\n"
        "and maintenance work delivered on the salon platform —\n"
        "much of it invisible from the user interface, but essential\n"
        "to the product's stability, scalability and growth.",
        size=14,
        color=INK,
    )

    # Numbers strip
    metrics = [
        ("186", "Commits"),
        ("75", "Pull Requests"),
        ("57", "DB Migrations"),
        ("7", "Months Active"),
    ]
    bx = Inches(5.2)
    by = Inches(6.0)
    bw = Inches(1.85)
    bh = Inches(1.0)
    gap = Inches(0.05)
    for i, (n, lbl) in enumerate(metrics):
        left = bx + (bw + gap) * i
        add_round_rect(s, left, by, bw, bh, LIGHT)
        add_text(s, left, by + Inches(0.1), bw, Inches(0.55), n, size=26, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_text(s, left, by + Inches(0.6), bw, Inches(0.35), lbl, size=10, color=MUTED, align=PP_ALIGN.CENTER)


def slide_agenda(prs):
    s = blank_slide(prs)
    add_page_header(s, "Contents", "What this report covers")

    items = [
        ("01", "Project overview", "Scope, team and engagement"),
        ("02", "Development timeline", "Month-by-month delivery"),
        ("03", "Git activity", "Visible evidence of continuous work"),
        ("04", "Backend & database", "Foundations under the UI"),
        ("05", "Features delivered", "What changed for users and the business"),
        ("06", "Maintenance & stability", "The invisible engineering layer"),
        ("07", "Metrics & impact", "Quantified output"),
        ("08", "Technical context", "How the pieces fit together"),
        ("09", "Conclusion", "Commitment going forward"),
    ]

    col_w = Inches(6.0)
    row_h = Inches(0.55)
    start_y = Inches(1.7)
    for i, (num, title, desc) in enumerate(items):
        col = i // 5
        row = i % 5
        left = Inches(0.6) + col * Inches(6.4)
        top = start_y + row * (row_h + Inches(0.15))
        add_round_rect(s, left, top, Inches(0.7), row_h, PRIMARY)
        add_text(s, left, top, Inches(0.7), row_h, num, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left + Inches(0.85), top + Inches(0.02), col_w, Inches(0.3), title, size=14, bold=True, color=NAVY)
        add_text(s, left + Inches(0.85), top + Inches(0.28), col_w, Inches(0.3), desc, size=10.5, color=MUTED)

    add_footer(s, "Page 2")


def slide_overview(prs):
    s = blank_slide(prs)
    add_page_header(s, "01 · Project Overview", "A platform under continuous engineering")

    body = (
        "The salon platform is a multi-tenant booking and management system serving "
        "customers, salon vendors, relationship managers and administrators. The "
        "backend, database, integrations and supporting infrastructure are all "
        "actively maintained — not as a one-time build, but as a living product."
    )
    add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1.1), body, size=14, color=INK)

    # Four overview cards
    cards = [
        (
            "Scope of work",
            [
                "Backend APIs (FastAPI)",
                "PostgreSQL via Supabase",
                "Payments (Razorpay)",
                "Storage & uploads",
                "Email & OTP services",
            ],
        ),
        (
            "Engagement",
            [
                "Period: Nov 2025 – May 2026",
                "Approx. 7 months active",
                "Continuous PR-based delivery",
                "Staging + production environments",
            ],
        ),
        (
            "Team responsibilities",
            [
                "Feature development",
                "Bug fixing & hotfixes",
                "Database migrations",
                "Deployment & operations",
                "Security & access control",
            ],
        ),
        (
            "User-facing surfaces",
            [
                "Customer mobile app",
                "Salon admin panel",
                "RM (relationship manager) tools",
                "Admin / internal panel",
            ],
        ),
    ]

    cx = Inches(0.6)
    cy = Inches(3.0)
    cw = Inches(3.0)
    ch = Inches(3.9)
    gap = Inches(0.13)
    for i, (title, items) in enumerate(cards):
        left = cx + (cw + gap) * i
        add_rect(s, left, cy, cw, ch, SOFT, line=LINE)
        add_rect(s, left, cy, cw, Inches(0.08), ACCENT)
        add_text(s, left + Inches(0.25), cy + Inches(0.2), cw - Inches(0.4), Inches(0.4), title, size=14, bold=True, color=NAVY)
        for j, it in enumerate(items):
            top = cy + Inches(0.7) + Inches(0.42) * j
            add_oval(s, left + Inches(0.28), top + Inches(0.13), Inches(0.08), Inches(0.08), ACCENT)
            add_text(s, left + Inches(0.45), top, cw - Inches(0.55), Inches(0.4), it, size=11, color=INK)

    add_footer(s, "Page 3")


def slide_timeline(prs):
    s = blank_slide(prs)
    add_page_header(s, "02 · Development Timeline", "Continuous delivery, month by month")

    add_text(
        s,
        Inches(0.6),
        Inches(1.55),
        Inches(12),
        Inches(0.5),
        "Each marker represents a phase of work. Commit volume varies by month, "
        "but engineering effort remained continuous throughout.",
        size=12,
        color=MUTED,
    )

    months = [
        ("Nov 2025", 38, "Foundation", ["Production schema", "Auth & RM flows", "Storage setup", "Career applications"]),
        ("Dec 2025", 4, "Stabilisation", ["Career APIs", "Document uploads", "Activity logs"]),
        ("Jan 2026", 22, "Auth & Ops", ["Token validity", "Profiler & locust", "CI pipeline", "Test coverage"]),
        ("Feb 2026", 30, "Compliance", ["Agreement uploads", "Payment reminders", "Business hours", "Email reliability"]),
        ("Mar 2026", 12, "Vendor", ["Outlet & GST", "OTP service", "Vendor approvals"]),
        ("Apr 2026", 51, "Features", ["Facilities", "Gender categories", "Reviews", "Favourites", "Discounts"]),
        ("May 2026", 28, "B2B", ["Products module", "Cart & checkout", "B2B pricing", "Subcategories", "CORS hardening"]),
    ]

    # Timeline axis
    ax_left = Inches(0.8)
    ax_right = Inches(12.7)
    ax_y = Inches(4.0)
    add_line(s, ax_left, ax_y, ax_right, ax_y, color=LINE, weight=2)

    n = len(months)
    span = ax_right - ax_left
    step = span // (n - 1)

    max_c = max(m[1] for m in months)
    for i, (label, count, phase, items) in enumerate(months):
        cx = ax_left + step * i
        # vertical bar above axis proportional to commits
        bar_h = Inches(0.5 + 1.8 * (count / max_c))
        bar_w = Inches(0.5)
        bar_left = cx - bar_w // 2
        bar_top = ax_y - bar_h
        add_rect(s, bar_left, bar_top, bar_w, bar_h, ACCENT)
        add_text(s, bar_left - Inches(0.3), bar_top - Inches(0.35), Inches(1.1), Inches(0.3), f"{count}", size=11, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

        # marker dot
        add_oval(s, cx - Inches(0.1), ax_y - Inches(0.1), Inches(0.2), Inches(0.2), PRIMARY)

        # month label
        add_text(s, cx - Inches(0.7), ax_y + Inches(0.1), Inches(1.4), Inches(0.3), label, size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # phase
        add_text(s, cx - Inches(0.7), ax_y + Inches(0.42), Inches(1.4), Inches(0.3), phase, size=9.5, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)

        # items - small text below
        text_lines = "\n".join("• " + it for it in items)
        add_text(s, cx - Inches(0.85), ax_y + Inches(0.8), Inches(1.7), Inches(2.2), text_lines, size=8.5, color=MUTED, align=PP_ALIGN.CENTER)

    # legend
    add_text(s, Inches(0.6), Inches(2.05), Inches(2), Inches(0.3), "Commit volume by month", size=10, bold=True, color=NAVY)
    add_rect(s, Inches(0.6), Inches(2.35), Inches(0.25), Inches(0.15), ACCENT)
    add_text(s, Inches(0.95), Inches(2.28), Inches(3), Inches(0.3), "Bar height = commits in that month", size=9.5, color=MUTED)

    add_footer(s, "Page 4")


def slide_git_activity(prs):
    s = blank_slide(prs)
    add_page_header(s, "03 · Git Activity & Contribution", "Each commit is a unit of engineering work")

    # Left explainer
    add_text(s, Inches(0.6), Inches(1.55), Inches(6.4), Inches(0.5), "What a commit represents", size=14, bold=True, color=NAVY)
    add_text(
        s,
        Inches(0.6),
        Inches(1.95),
        Inches(6.4),
        Inches(2.4),
        "Software maintenance often happens behind the scenes. While the user "
        "interface may appear similar between releases, continuous engineering "
        "work is required to improve stability, scalability, security and "
        "performance.\n\n"
        "Each commit on the right may include any combination of:",
        size=12,
        color=INK,
    )

    bullets = [
        "Multiple file updates",
        "Backend logic improvements",
        "Database schema changes",
        "New feature code",
        "Bug fixes and hotfixes",
        "Refactoring of existing code",
        "Performance optimisations",
        "Security & access-control changes",
    ]
    for i, b in enumerate(bullets):
        col = i // 4
        row = i % 4
        left = Inches(0.6) + col * Inches(3.2)
        top = Inches(4.05) + row * Inches(0.32)
        add_oval(s, left, top + Inches(0.1), Inches(0.1), Inches(0.1), ACCENT)
        add_text(s, left + Inches(0.2), top, Inches(3.0), Inches(0.3), b, size=11, color=INK)

    # Right: commit heatmap (7 rows x 28 cols approx) using real monthly distribution
    grid_left = Inches(7.4)
    grid_top = Inches(1.7)
    cell = Inches(0.18)
    gap = Inches(0.04)
    cols = 26
    rows = 7

    monthly = [38, 4, 22, 30, 12, 51, 28]
    total_cells = cols * rows
    # Distribute commits roughly per month over equal column blocks
    block = cols // len(monthly)
    cell_count = [[0] * cols for _ in range(rows)]

    import random
    rng = random.Random(7)

    for mi, count in enumerate(monthly):
        c_start = mi * block
        c_end = min(c_start + block, cols)
        cells_in_block = (c_end - c_start) * rows
        remaining = count
        for _ in range(count):
            r = rng.randint(0, rows - 1)
            c = rng.randint(c_start, c_end - 1)
            cell_count[r][c] += 1

    max_v = max(max(row) for row in cell_count) or 1
    add_text(s, Inches(7.4), Inches(1.4), Inches(5.4), Inches(0.3), "Commit heatmap (illustrative)", size=11, bold=True, color=NAVY)

    for r in range(rows):
        for c in range(cols):
            v = cell_count[r][c]
            if v == 0:
                color = LIGHT
            else:
                ratio = v / max_v
                if ratio < 0.34:
                    color = RGBColor(0xB6, 0xD4, 0xF1)
                elif ratio < 0.67:
                    color = ACCENT
                else:
                    color = PRIMARY
            x = grid_left + c * (cell + gap)
            y = grid_top + r * (cell + gap)
            add_rect(s, x, y, cell, cell, color)

    # heatmap legend
    leg_y = grid_top + rows * (cell + gap) + Inches(0.15)
    add_text(s, Inches(7.4), leg_y, Inches(2), Inches(0.3), "Less", size=9, color=MUTED)
    for i, col in enumerate([LIGHT, RGBColor(0xB6, 0xD4, 0xF1), ACCENT, PRIMARY]):
        add_rect(s, Inches(7.85) + i * Inches(0.32), leg_y + Inches(0.04), Inches(0.25), Inches(0.18), col)
    add_text(s, Inches(9.2), leg_y, Inches(2), Inches(0.3), "More", size=9, color=MUTED)

    # Contributor mini chart
    contrib_y = Inches(4.7)
    add_text(s, Inches(7.4), contrib_y, Inches(5.4), Inches(0.3), "Contributors (commits)", size=11, bold=True, color=NAVY)
    contribs = [("ve-sca-via", 85), ("Safdar Ali Niazi", 40), ("safdaraliniazi", 37), ("Your Name", 31)]
    max_cv = max(c for _, c in contribs)
    bar_left = Inches(8.9)
    bar_max_w = Inches(3.7)
    for i, (name, count) in enumerate(contribs):
        y = contrib_y + Inches(0.35) + i * Inches(0.4)
        add_text(s, Inches(7.4), y, Inches(1.45), Inches(0.3), name, size=10, color=INK)
        w = Emu(int(bar_max_w * count / max_cv))
        add_rect(s, bar_left, y + Inches(0.05), w, Inches(0.22), ACCENT)
        add_text(s, bar_left + w + Inches(0.05), y, Inches(0.7), Inches(0.3), str(count), size=10, color=PRIMARY, bold=True)

    add_footer(s, "Page 5")


def slide_backend_database(prs):
    s = blank_slide(prs)
    add_page_header(s, "04 · Backend & Database Work", "Foundations that sit beneath every screen")

    add_text(
        s,
        Inches(0.6),
        Inches(1.6),
        Inches(12),
        Inches(0.6),
        "Most of the engineering effort lives below the UI. These are the categories "
        "of backend and database work delivered during the engagement.",
        size=12,
        color=MUTED,
    )

    sections = [
        ("Database structure", [
            "57 versioned migrations applied",
            "New tables: products, orders, carts, OTPs, activity logs, email logs, careers",
            "Service subcategories, gender categories, facilities, business hours",
        ]),
        ("APIs & business logic", [
            "16 API modules covering vendors, salons, RM, admin, payments",
            "Booking, products, cart and checkout flows",
            "Review system, favourites, discounts, B2B pricing",
        ]),
        ("Authentication & security", [
            "Custom JWT + Supabase auth flow consolidated",
            "Token validity timestamps & forced logout-all",
            "OTP service (MessageCentral) with rate limiting and lockouts",
            "Password reset / forgot-password flows",
            "Role-based access (admin / RM / vendor / customer)",
            "Storage bucket policies & private agreement uploads",
        ]),
        ("Reliability & ops", [
            "Database client singleton & refresh logic",
            "Email delivery hardening + email_logs table",
            "Activity logging for admin visibility",
            "CORS preflight fixes & deployment configuration",
            "Locust load profile & profiler integration",
            "CI pipeline for safer deployments",
        ]),
    ]

    cx = Inches(0.6)
    cy = Inches(2.4)
    cw = Inches(6.0)
    ch = Inches(2.3)
    gap = Inches(0.15)

    for i, (title, items) in enumerate(sections):
        col = i % 2
        row = i // 2
        left = cx + col * (cw + gap)
        top = cy + row * (ch + gap)
        add_rect(s, left, top, cw, ch, WHITE, line=LINE)
        add_rect(s, left, top, Inches(0.12), ch, ACCENT)
        add_text(s, left + Inches(0.3), top + Inches(0.15), cw - Inches(0.4), Inches(0.4), title, size=14, bold=True, color=NAVY)
        for j, it in enumerate(items):
            ty = top + Inches(0.6) + Inches(0.3) * j
            add_oval(s, left + Inches(0.35), ty + Inches(0.1), Inches(0.08), Inches(0.08), ACCENT)
            add_text(s, left + Inches(0.55), ty, cw - Inches(0.7), Inches(0.3), it, size=10.5, color=INK)

    add_footer(s, "Page 6")


def slide_features(prs):
    s = blank_slide(prs)
    add_page_header(s, "05 · New Features Delivered", "Capabilities added during this period")

    features = [
        ("B2B Products & Wholesale", "Salons could only sell services. No path to sell or buy products.", "Full products module with cart, checkout, payment, and B2B (wholesale) pricing for regular buyers.", "Database tables, APIs, payment flow, B2B price logic, admin order management."),
        ("Hierarchical Service Categories", "Flat category list limited how services could be organised.", "Subcategories with vendor & RM support — better browsing and admin control.", "service_subcategories table, CRUD APIs, vendor & RM integration."),
        ("Reviews & Favourites", "No social proof, no way for users to save salons.", "Customers can review salons and favourite ones they like.", "Reviews schema & APIs, favourites table, customer endpoints."),
        ("Vendor Onboarding (facilities, GST, outlets)", "Vendor data captured only minimal information.", "Richer vendor profiles for compliance and discoverability.", "JSONB facilities, gender categories, outlet & GST fields, approval logic."),
        ("OTP Login (MessageCentral)", "Email-only login restricted some user segments.", "Phone-based OTP login with rate limiting and lockouts.", "otp_service, otp_attempts table, send/verify endpoints, security hardening."),
        ("Career Applications", "No structured way to receive applications.", "End-to-end career application capture with document uploads.", "Schema, status flow, Cloudinary integration, status update API."),
    ]

    cw = Inches(4.05)
    ch = Inches(2.65)
    cx = Inches(0.6)
    cy = Inches(1.55)
    gap_x = Inches(0.15)
    gap_y = Inches(0.15)

    for i, (name, problem, impact, tech) in enumerate(features):
        col = i % 3
        row = i // 3
        left = cx + col * (cw + gap_x)
        top = cy + row * (ch + gap_y)
        add_rect(s, left, top, cw, ch, SOFT, line=LINE)
        add_rect(s, left, top, cw, Inches(0.45), PRIMARY)
        add_text(s, left + Inches(0.2), top + Inches(0.05), cw - Inches(0.3), Inches(0.4), name, size=12, bold=True, color=WHITE)

        ty = top + Inches(0.55)
        add_text(s, left + Inches(0.2), ty, cw - Inches(0.3), Inches(0.25), "PROBLEM", size=8, bold=True, color=AMBER)
        add_text(s, left + Inches(0.2), ty + Inches(0.22), cw - Inches(0.3), Inches(0.5), problem, size=9.5, color=INK)
        ty = top + Inches(1.2)
        add_text(s, left + Inches(0.2), ty, cw - Inches(0.3), Inches(0.25), "DELIVERED", size=8, bold=True, color=GREEN)
        add_text(s, left + Inches(0.2), ty + Inches(0.22), cw - Inches(0.3), Inches(0.55), impact, size=9.5, color=INK)
        ty = top + Inches(1.95)
        add_text(s, left + Inches(0.2), ty, cw - Inches(0.3), Inches(0.25), "TECHNICAL WORK", size=8, bold=True, color=ACCENT)
        add_text(s, left + Inches(0.2), ty + Inches(0.22), cw - Inches(0.3), Inches(0.55), tech, size=9.5, color=INK)

    add_footer(s, "Page 7")


def slide_maintenance(prs):
    s = blank_slide(prs)
    add_page_header(s, "06 · Maintenance & Stability", "The work users don't see")

    # Analogy card
    add_round_rect(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(1.1), LIGHT)
    add_text(s, Inches(0.85), Inches(1.65), Inches(11.6), Inches(0.45), "An analogy", size=11, bold=True, color=PRIMARY)
    add_text(
        s,
        Inches(0.85),
        Inches(1.95),
        Inches(11.6),
        Inches(0.7),
        "Like maintaining a car engine, many important improvements happen internally and "
        "may not change the exterior appearance. The vehicle still looks the same — but it "
        "starts more reliably, runs more efficiently, and is safer to drive.",
        size=12,
        color=INK,
    )

    categories = [
        ("Bug fixing", "22 explicit fix commits, plus inline corrections in feature PRs."),
        ("Downtime prevention", "Email reliability, payment idempotency, retry-safe order flows."),
        ("Performance optimisation", "LRU cache for Razorpay credentials, geocoding cache, query tuning."),
        ("Dependency updates", "Python runtime pinned, Pillow/SDK updates, deployment config cleanup."),
        ("Security patches", "Token rotation, storage policies, OTP rate limits, CORS hardening."),
        ("Refactoring", "Auth flow consolidation, duplicate code removal, sync supabase client."),
        ("Scalability", "Locust load tests, profiler instrumentation, async-safe payment flow."),
        ("Monitoring & logs", "Activity logs, email logs, structured rich logging."),
    ]

    cx = Inches(0.6)
    cy = Inches(2.95)
    cw = Inches(3.0)
    ch = Inches(2.05)
    gap = Inches(0.13)
    for i, (title, desc) in enumerate(categories):
        col = i % 4
        row = i // 4
        left = cx + col * (cw + gap)
        top = cy + row * (ch + gap)
        add_rect(s, left, top, cw, ch, WHITE, line=LINE)
        add_oval(s, left + Inches(0.25), top + Inches(0.25), Inches(0.45), Inches(0.45), LIGHT)
        add_text(s, left + Inches(0.25), top + Inches(0.25), Inches(0.45), Inches(0.45), str(i + 1), size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left + Inches(0.85), top + Inches(0.25), cw - Inches(1.0), Inches(0.4), title, size=12.5, bold=True, color=NAVY)
        add_text(s, left + Inches(0.25), top + Inches(0.85), cw - Inches(0.4), Inches(1.1), desc, size=10, color=MUTED)

    add_footer(s, "Page 8")


def slide_metrics(prs):
    s = blank_slide(prs)
    add_page_header(s, "07 · Metrics & Impact", "Quantified output, drawn directly from the repository")

    # KPI cards row
    kpis = [
        ("186", "Commits", "Individual code changes shipped"),
        ("75", "Pull requests merged", "Reviewed change sets"),
        ("57", "Database migrations", "Versioned schema changes"),
        ("16", "API modules", "Backend feature areas"),
        ("24", "Service modules", "Business-logic layers"),
        ("~115K", "Lines added", "Code authored across the period"),
    ]
    cx = Inches(0.6)
    cy = Inches(1.55)
    cw = Inches(2.0)
    ch = Inches(1.45)
    gap = Inches(0.07)
    for i, (n, lbl, sub) in enumerate(kpis):
        left = cx + (cw + gap) * i
        add_rect(s, left, cy, cw, ch, PRIMARY)
        add_text(s, left, cy + Inches(0.15), cw, Inches(0.6), n, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, left, cy + Inches(0.75), cw, Inches(0.3), lbl, size=10.5, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
        add_text(s, left, cy + Inches(1.05), cw, Inches(0.35), sub, size=9, color=LIGHT, align=PP_ALIGN.CENTER)

    # Chart: commits per month
    chart_left = Inches(0.6)
    chart_top = Inches(3.25)
    chart_w = Inches(7.8)
    chart_h = Inches(3.7)
    add_rect(s, chart_left, chart_top, chart_w, chart_h, SOFT, line=LINE)
    add_text(s, chart_left + Inches(0.25), chart_top + Inches(0.15), chart_w - Inches(0.5), Inches(0.3), "Commits per month", size=12, bold=True, color=NAVY)

    chart_data = CategoryChartData()
    chart_data.categories = ["Nov '25", "Dec '25", "Jan '26", "Feb '26", "Mar '26", "Apr '26", "May '26"]
    chart_data.add_series("Commits", (38, 4, 22, 30, 12, 51, 28))
    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        chart_left + Inches(0.2),
        chart_top + Inches(0.55),
        chart_w - Inches(0.4),
        chart_h - Inches(0.75),
        chart_data,
    ).chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(10)
    plot.data_labels.font.bold = True
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    series = plot.series[0]
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT
    series.format.line.fill.background()

    # Donut: work mix (approx based on commit message analysis)
    donut_left = Inches(8.6)
    donut_top = Inches(3.25)
    donut_w = Inches(4.1)
    donut_h = Inches(3.7)
    add_rect(s, donut_left, donut_top, donut_w, donut_h, SOFT, line=LINE)
    add_text(s, donut_left + Inches(0.25), donut_top + Inches(0.15), donut_w - Inches(0.5), Inches(0.3), "Work distribution (commits)", size=12, bold=True, color=NAVY)

    donut_data = CategoryChartData()
    donut_data.categories = ["Features", "Fixes", "Refactor / maint", "Infra / deploy", "Merges & releases"]
    donut_data.add_series("Mix", (38, 32, 22, 19, 75))
    donut = s.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        donut_left + Inches(0.2),
        donut_top + Inches(0.55),
        donut_w - Inches(0.4),
        donut_h - Inches(0.75),
        donut_data,
    ).chart
    donut.has_title = False
    donut.has_legend = True
    donut.legend.position = XL_LEGEND_POSITION.RIGHT
    donut.legend.include_in_layout = False
    donut.legend.font.size = Pt(9)

    add_footer(s, "Page 9")


def slide_technical_context(prs):
    s = blank_slide(prs)
    add_page_header(s, "08 · Technical Context, Simplified", "How UI, backend and data fit together")

    # Architecture diagram - 4 horizontal layers
    layers = [
        ("User Interfaces", "Mobile app · Salon admin · RM tools · Admin panel", LIGHT, NAVY),
        ("API Layer (FastAPI)", "Authentication · Bookings · Products · Payments · RM · Vendors · Admin", RGBColor(0xCF, 0xE3, 0xF8), PRIMARY),
        ("Services & Business Logic", "Orders · Pricing · OTP · Email · Activity logs · Geocoding · Storage", RGBColor(0xA8, 0xC9, 0xEA), PRIMARY),
        ("Data & Infrastructure", "PostgreSQL (Supabase) · Storage buckets · Razorpay · Cloudinary · MessageCentral", PRIMARY, WHITE),
    ]

    left = Inches(0.8)
    top = Inches(1.65)
    width = Inches(8.0)
    layer_h = Inches(1.05)
    gap = Inches(0.15)

    for i, (title, sub, fill, ink) in enumerate(layers):
        y = top + i * (layer_h + gap)
        add_round_rect(s, left, y, width, layer_h, fill)
        add_text(s, left + Inches(0.3), y + Inches(0.12), width - Inches(0.5), Inches(0.4), title, size=14, bold=True, color=ink)
        add_text(s, left + Inches(0.3), y + Inches(0.5), width - Inches(0.5), Inches(0.5), sub, size=10.5, color=ink)
        # Arrows between layers
        if i < len(layers) - 1:
            arrow_x = left + width // 2
            arrow_top = y + layer_h
            arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, arrow_x - Inches(0.15), arrow_top - Inches(0.02), Inches(0.3), gap + Inches(0.04))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()

    # Right-side explainer
    rx = Inches(9.2)
    ry = Inches(1.65)
    rw = Inches(3.7)
    add_rect(s, rx, ry, rw, Inches(5.3), SOFT, line=LINE)
    add_rect(s, rx, ry, rw, Inches(0.08), ACCENT)
    add_text(s, rx + Inches(0.25), ry + Inches(0.2), rw - Inches(0.5), Inches(0.4), "Why this matters", size=13, bold=True, color=NAVY)

    notes = [
        ("UI is the tip of the iceberg.", "Every screen depends on dozens of API endpoints, services and database tables."),
        ("Maintenance is structural.", "Schema changes, auth hardening and payment integrity have no visible UI — but the platform cannot function safely without them."),
        ("Data flows are bidirectional.", "A simple action (e.g. booking) touches auth, services, salons, payments, storage and logs."),
        ("Changes compound.", "Each improvement to a lower layer benefits every feature built above it."),
    ]
    ty = ry + Inches(0.65)
    for h, body in notes:
        add_text(s, rx + Inches(0.25), ty, rw - Inches(0.5), Inches(0.3), h, size=11, bold=True, color=PRIMARY)
        add_text(s, rx + Inches(0.25), ty + Inches(0.3), rw - Inches(0.5), Inches(0.9), body, size=10, color=INK)
        ty += Inches(1.15)

    add_footer(s, "Page 10")


def slide_conclusion(prs):
    s = blank_slide(prs)

    # Big banner
    add_rect(s, 0, 0, SLIDE_W, Inches(2.2), NAVY)
    add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4), "09 · CONCLUSION", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8), "Continuous engineering, transparent reporting.", size=32, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.4), "A summary of the commitments going forward.", size=14, color=LIGHT)

    # Three commitment cards
    cards = [
        ("Stability first", "Continued investment in monitoring, logging, automated tests and load profiling so issues are detected and resolved early — often before users encounter them."),
        ("Scalability & performance", "Ongoing refinement of database queries, caching, and async paths so the platform absorbs growth in traffic, vendors and product orders without degradation."),
        ("Transparent collaboration", "Every change is reviewed, versioned and deployable. This report format can be repeated on a regular cadence so progress remains visible and measurable."),
    ]

    cx = Inches(0.6)
    cy = Inches(2.6)
    cw = Inches(4.05)
    ch = Inches(3.4)
    gap = Inches(0.13)
    for i, (title, body) in enumerate(cards):
        left = cx + (cw + gap) * i
        add_rect(s, left, cy, cw, ch, SOFT, line=LINE)
        add_rect(s, left, cy, cw, Inches(0.1), ACCENT)
        add_oval(s, left + Inches(0.3), cy + Inches(0.4), Inches(0.7), Inches(0.7), LIGHT)
        add_text(s, left + Inches(0.3), cy + Inches(0.4), Inches(0.7), Inches(0.7), str(i + 1), size=20, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left + Inches(0.3), cy + Inches(1.25), cw - Inches(0.5), Inches(0.5), title, size=16, bold=True, color=NAVY)
        add_text(s, left + Inches(0.3), cy + Inches(1.85), cw - Inches(0.5), Inches(1.5), body, size=11.5, color=INK)

    # Closing strip
    add_rect(s, 0, Inches(6.25), SLIDE_W, Inches(1.25), LIGHT)
    add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.4), "Thank you.", size=20, bold=True, color=NAVY)
    add_text(
        s,
        Inches(0.6),
        Inches(6.8),
        Inches(12),
        Inches(0.5),
        "We remain committed to maintaining and improving the platform — both the visible product and the engineering foundation it depends on.",
        size=12,
        color=INK,
    )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_agenda(prs)
    slide_overview(prs)
    slide_timeline(prs)
    slide_git_activity(prs)
    slide_backend_database(prs)
    slide_features(prs)
    slide_maintenance(prs)
    slide_metrics(prs)
    slide_technical_context(prs)
    slide_conclusion(prs)

    out = Path("Development_Progress_Report.pptx")
    prs.save(out)
    print(f"Wrote {out.resolve()} ({out.stat().st_size / 1024:.1f} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
